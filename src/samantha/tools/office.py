"""Safe local Office automation for Samantha."""
# ruff: noqa: E501

from __future__ import annotations

from pathlib import Path
from typing import Any

from samantha.core.registry import ToolRegistry
from samantha.core.types import ToolResult
from samantha.tools._stubs import BaseTool, ToolSpec


@ToolRegistry.register("office")
class OfficeTool(BaseTool):
    """Inspect and update Excel, Word, and PowerPoint files."""

    tool_id = "office"

    @property
    def spec(self) -> ToolSpec:
        return ToolSpec(
            name="office",
            description=(
                "Work with local Excel, Word, and PowerPoint files. "
                "Excel actions include inspect_workbook, read_range, write_cells, "
                "set_formulas, and create_dashboard. Write operations require confirmation."
            ),
            parameters={
                "type": "object",
                "properties": {
                    "action": {"type": "string"},
                    "path": {"type": "string"},
                    "sheet": {"type": "string"},
                    "range": {"type": "string"},
                    "values": {"type": "array"},
                    "formulas": {"type": "object"},
                    "dashboard_name": {"type": "string"},
                },
                "required": ["action", "path"],
                "additionalProperties": False,
            },
            category="office",
            requires_confirmation=True,
            timeout_seconds=60.0,
            required_capabilities=["file:write"],
        )

    def execute(self, **params: Any) -> ToolResult:
        action = str(params.get("action", ""))
        path = Path(str(params.get("path", ""))).expanduser().resolve()
        if not path.exists():
            return ToolResult(
                tool_name=self.tool_id, content=f"File not found: {path}", success=False
            )
        try:
            if action == "inspect_workbook":
                return self._inspect_workbook(path)
            if action == "read_range":
                return self._read_range(path, params)
            if action in {"write_cells", "set_formulas", "create_dashboard"}:
                return self._write_excel(path, action, params)
            if action == "inspect_word":
                from docx import Document

                text = "\n".join(
                    p.text for p in Document(path).paragraphs if p.text.strip()
                )
                return ToolResult(
                    tool_name=self.tool_id, content=text[:50000], success=True
                )
            if action == "inspect_powerpoint":
                from pptx import Presentation

                prs = Presentation(path)
                lines = []
                for i, slide in enumerate(prs.slides, 1):
                    text = " ".join(
                        s.text for s in slide.shapes if hasattr(s, "text") and s.text
                    )
                    lines.append(f"Slide {i}: {text}")
                return ToolResult(
                    tool_name=self.tool_id,
                    content="\n".join(lines)[:50000],
                    success=True,
                )
            return ToolResult(
                tool_name=self.tool_id,
                content=f"Unsupported Office action: {action}",
                success=False,
            )
        except ImportError as exc:
            return ToolResult(
                tool_name=self.tool_id,
                content=f"Missing Office dependency: {exc}",
                success=False,
            )
        except Exception as exc:
            return ToolResult(
                tool_name=self.tool_id,
                content=f"Office operation failed: {exc}",
                success=False,
            )

    @staticmethod
    def _workbook(path: Path):
        from openpyxl import load_workbook

        return load_workbook(
            path, data_only=False, keep_vba=path.suffix.lower() == ".xlsm"
        )

    def _inspect_workbook(self, path: Path) -> ToolResult:
        wb = self._workbook(path)
        info = [
            f"{ws.title}: {ws.max_row} rows x {ws.max_column} columns"
            for ws in wb.worksheets
        ]
        return ToolResult(tool_name=self.tool_id, content="\n".join(info), success=True)

    def _read_range(self, path: Path, params: dict[str, Any]) -> ToolResult:
        from openpyxl.utils.cell import range_boundaries

        wb = self._workbook(path)
        ws = wb[str(params.get("sheet", wb.sheetnames[0]))]
        rng = str(params.get("range", "A1:Z50"))
        min_col, min_row, max_col, max_row = range_boundaries(rng)
        if max_row - min_row > 5000:
            return ToolResult(
                tool_name=self.tool_id, content="Range is too large.", success=False
            )
        rows = [
            [ws.cell(r, c).value for c in range(min_col, max_col + 1)]
            for r in range(min_row, max_row + 1)
        ]
        return ToolResult(tool_name=self.tool_id, content=repr(rows), success=True)

    def _write_excel(
        self, path: Path, action: str, params: dict[str, Any]
    ) -> ToolResult:
        import shutil
        from datetime import datetime

        wb = self._workbook(path)
        ws = wb[str(params.get("sheet", wb.sheetnames[0]))]
        backup = path.with_name(
            f"{path.stem}.samantha-backup-{datetime.now():%Y%m%d-%H%M%S}{path.suffix}"
        )
        shutil.copy2(path, backup)
        if action == "write_cells":
            from openpyxl.utils.cell import coordinate_to_tuple

            values = params.get("values")
            row, col = coordinate_to_tuple(str(params.get("range", "A1")).split(":")[0])
            for r, values_row in enumerate(values or []):
                for c, value in enumerate(values_row):
                    ws.cell(row + r, col + c).value = value
        elif action == "set_formulas":
            for cell, formula in (params.get("formulas") or {}).items():
                value = str(formula)
                ws[str(cell)].value = value if value.startswith("=") else "=" + value
        elif action == "create_dashboard":
            name = str(params.get("dashboard_name", "Dashboard"))[:31]
            if name in wb.sheetnames:
                del wb[name]
            dash = wb.create_sheet(name, 0)
            dash["A1"] = "Samantha Dashboard"
            dash["A3"] = "Sheet"
            dash["B3"] = "Rows"
            dash["C3"] = "Columns"
            for i, source in enumerate(wb.worksheets[1:], 4):
                dash.cell(i, 1).value = source.title
                dash.cell(i, 2).value = source.max_row
                dash.cell(i, 3).value = source.max_column
        wb.save(path)
        return ToolResult(
            tool_name=self.tool_id,
            content=f"Updated {path}. Backup: {backup}",
            success=True,
        )


__all__ = ["OfficeTool"]

"""Safe Microsoft Office-style document automation for Samantha.

Supports local XLSX workbooks plus DOCX and PPTX creation. Mutating workbook
operations create a timestamped backup before changing the source file.
"""

from __future__ import annotations

import json
import shutil
from datetime import datetime
from pathlib import Path
from typing import Any

from samantha.core.registry import ToolRegistry
from samantha.core.types import ToolResult
from samantha.tools._stubs import BaseTool, ToolSpec

_MAX_FILE_BYTES = 100 * 1024 * 1024
_ALLOWED_EXTENSIONS = {".xlsx", ".xlsm", ".docx", ".pptx"}


def _result(content: str, success: bool = True, **metadata: Any) -> ToolResult:
    return ToolResult(
        tool_name="office_control",
        content=content,
        success=success,
        metadata=metadata or None,
    )


def _validate_path(raw: str, *, extension: str | None = None) -> Path:
    path = Path(raw).expanduser().resolve()
    if len(str(path)) > 4096:
        raise ValueError("Path is too long.")
    if extension and path.suffix.lower() != extension:
        raise ValueError(f"Expected a {extension} file.")
    if path.suffix.lower() not in _ALLOWED_EXTENSIONS:
        raise ValueError("Unsupported Office file type.")
    if path.exists() and path.stat().st_size > _MAX_FILE_BYTES:
        raise ValueError("File is larger than the 100 MB safety limit.")
    return path


def _backup(path: Path) -> Path:
    stamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    backup = path.with_name(f"{path.stem}.samantha-backup-{stamp}{path.suffix}")
    shutil.copy2(path, backup)
    return backup


@ToolRegistry.register("office_control")
class OfficeControlTool(BaseTool):
    """Inspect and safely modify common Microsoft Office file formats."""

    tool_id = "office_control"

    @property
    def spec(self) -> ToolSpec:
        return ToolSpec(
            name="office_control",
            description=(
                "Work with local Excel, Word and PowerPoint files. "
                "Inspect XLSX sheets/formulas, read and write ranges, repair "
                "formulas, create dashboard sheets, and create DOCX/PPTX files. "
                "Never execute macros. Mutating actions require confirmation."
            ),
            parameters={
                "type": "object",
                "properties": {
                    "action": {
                        "type": "string",
                        "enum": [
                            "inspect_workbook",
                            "read_range",
                            "write_range",
                            "set_formula",
                            "create_dashboard_sheet",
                            "create_docx",
                            "create_pptx",
                        ],
                    },
                    "path": {"type": "string"},
                    "sheet": {"type": "string"},
                    "range": {"type": "string"},
                    "values": {
                        "type": "array",
                        "description": "2D array for an Excel range.",
                    },
                    "formula": {"type": "string"},
                    "output_path": {"type": "string"},
                    "title": {"type": "string"},
                    "content": {"type": "string"},
                    "slides": {"type": "array"},
                },
                "required": ["action"],
                "additionalProperties": False,
            },
            category="office",
            requires_confirmation=True,
            required_capabilities=["file:write"],
            timeout_seconds=60.0,
        )

    def execute(self, **params: Any) -> ToolResult:
        action = str(params.get("action", "")).strip()
        try:
            if action == "inspect_workbook":
                return self._inspect(params)
            if action == "read_range":
                return self._read_range(params)
            if action == "write_range":
                return self._write_range(params)
            if action == "set_formula":
                return self._set_formula(params)
            if action == "create_dashboard_sheet":
                return self._dashboard(params)
            if action == "create_docx":
                return self._docx(params)
            if action == "create_pptx":
                return self._pptx(params)
            return _result(f"Unsupported Office action: {action!r}", False)
        except Exception as exc:
            return _result(f"Office operation failed: {exc}", False)

    @staticmethod
    def _load_workbook(path: Path) -> Any:
        from openpyxl import load_workbook

        # read_only is intentionally false because formulas and styles need inspection.
        return load_workbook(
            path, data_only=False, keep_vba=path.suffix.lower() == ".xlsm"
        )

    def _inspect(self, p: dict[str, Any]) -> ToolResult:
        path = _validate_path(str(p.get("path", "")), extension=".xlsx")
        wb = self._load_workbook(path)
        sheets = []
        for ws in wb.worksheets:
            formula_count = 0
            nonempty = 0
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        nonempty += 1
                        if isinstance(cell.value, str) and cell.value.startswith("="):
                            formula_count += 1
            sheets.append(
                {
                    "name": ws.title,
                    "rows": ws.max_row,
                    "columns": ws.max_column,
                    "nonempty_cells": nonempty,
                    "formula_cells": formula_count,
                    "tables": list(ws.tables.keys()),
                }
            )
        wb.close()
        return _result(
            json.dumps({"path": str(path), "sheets": sheets}, indent=2),
            True,
            path=str(path),
            sheets=sheets,
        )

    def _read_range(self, p: dict[str, Any]) -> ToolResult:
        path = _validate_path(str(p.get("path", "")), extension=".xlsx")
        sheet = str(p.get("sheet", ""))
        cell_range = str(p.get("range", ""))
        wb = self._load_workbook(path)
        if sheet not in wb.sheetnames:
            wb.close()
            return _result(f"Sheet not found: {sheet}", False)
        ws = wb[sheet]
        values = [[cell.value for cell in row] for row in ws[cell_range]]
        wb.close()
        return _result(json.dumps(values, default=str), True, values=values)

    def _write_range(self, p: dict[str, Any]) -> ToolResult:
        path = _validate_path(str(p.get("path", "")), extension=".xlsx")
        sheet = str(p.get("sheet", ""))
        cell_range = str(p.get("range", ""))
        values = p.get("values")
        if not isinstance(values, list) or not values:
            return _result("values must be a non-empty 2D array.", False)
        if not path.exists():
            return _result(f"Workbook does not exist: {path}", False)
        backup = _backup(path)
        wb = self._load_workbook(path)
        if sheet not in wb.sheetnames:
            wb.close()
            return _result(f"Sheet not found: {sheet}", False)
        ws = wb[sheet]
        target = ws[cell_range]
        rows = target if isinstance(target, tuple) else ((target,),)
        for r_idx, row in enumerate(values):
            if not isinstance(row, list):
                row = [row]
            for c_idx, value in enumerate(row):
                if r_idx >= len(rows) or c_idx >= len(rows[r_idx]):
                    break
                rows[r_idx][c_idx].value = value
        wb.save(path)
        wb.close()
        return _result(
            f"Updated {sheet}!{cell_range} in {path.name}.", True, backup=str(backup)
        )

    def _set_formula(self, p: dict[str, Any]) -> ToolResult:
        path = _validate_path(str(p.get("path", "")), extension=".xlsx")
        sheet = str(p.get("sheet", ""))
        cell = str(p.get("range", ""))
        formula = str(p.get("formula", "")).strip()
        if not formula.startswith("="):
            formula = "=" + formula
        if not formula or len(formula) > 10000:
            return _result("Invalid formula.", False)
        backup = _backup(path)
        wb = self._load_workbook(path)
        if sheet not in wb.sheetnames:
            wb.close()
            return _result(f"Sheet not found: {sheet}", False)
        wb[sheet][cell] = formula
        wb.save(path)
        wb.close()
        return _result(f"Formula written to {sheet}!{cell}.", True, backup=str(backup))

    def _dashboard(self, p: dict[str, Any]) -> ToolResult:
        path = _validate_path(str(p.get("path", "")), extension=".xlsx")
        source_sheet = str(p.get("sheet", ""))
        title = (
            str(p.get("title", "Samantha Dashboard")).strip() or "Samantha Dashboard"
        )
        backup = _backup(path)
        wb = self._load_workbook(path)
        if source_sheet not in wb.sheetnames:
            wb.close()
            return _result(f"Sheet not found: {source_sheet}", False)
        if title in wb.sheetnames:
            del wb[title]
        dashboard = wb.create_sheet(title)
        source = wb[source_sheet]
        dashboard["A1"] = title
        dashboard["A3"] = "Source sheet"
        dashboard["B3"] = source_sheet
        dashboard["A4"] = "Rows"
        dashboard["B4"] = source.max_row
        dashboard["A5"] = "Columns"
        dashboard["B5"] = source.max_column
        dashboard["A7"] = "Top-left preview"
        max_rows = min(source.max_row, 15)
        max_cols = min(source.max_column, 10)
        for r in range(1, max_rows + 1):
            for c in range(1, max_cols + 1):
                dashboard.cell(row=8 + r, column=c).value = source.cell(
                    row=r, column=c
                ).value
        for column in dashboard.columns:
            width = min(max(max(len(str(c.value or "")) for c in column) + 2, 10), 40)
            dashboard.column_dimensions[column[0].column_letter].width = width
        wb.save(path)
        wb.close()
        return _result(
            f"Created '{title}' dashboard sheet in {path.name}.",
            True,
            backup=str(backup),
            sheet=title,
        )

    def _docx(self, p: dict[str, Any]) -> ToolResult:
        output = Path(str(p.get("output_path", ""))).expanduser().resolve()
        if output.suffix.lower() != ".docx":
            output = output.with_suffix(".docx")
        from docx import Document

        doc = Document()
        title = str(p.get("title", "Samantha Document"))
        doc.add_heading(title, level=1)
        content = str(p.get("content", ""))
        for paragraph in content.split("\n"):
            doc.add_paragraph(paragraph)
        output.parent.mkdir(parents=True, exist_ok=True)
        doc.save(output)
        return _result(f"Created Word document: {output}", True, path=str(output))

    def _pptx(self, p: dict[str, Any]) -> ToolResult:
        output = Path(str(p.get("output_path", ""))).expanduser().resolve()
        if output.suffix.lower() != ".pptx":
            output = output.with_suffix(".pptx")
        from pptx import Presentation

        prs = Presentation()
        title = str(p.get("title", "Samantha Presentation"))
        slides = p.get("slides") or [
            {"title": title, "body": str(p.get("content", ""))}
        ]
        for index, item in enumerate(slides):
            layout = prs.slide_layouts[1 if len(prs.slide_layouts) > 1 else 0]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = str(item.get("title", f"Slide {index + 1}"))
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = str(item.get("body", ""))
        output.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output)
        return _result(
            f"Created PowerPoint presentation: {output}", True, path=str(output)
        )


__all__ = ["OfficeControlTool"]
